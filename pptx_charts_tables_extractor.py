from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_charts(ppt_file_path):
    # 파워포인트 프레젠테이션 파일 로드
    presentation = Presentation(ppt_file_path)

    # 차트 데이터를 저장할 리스트
    charts_data = [[] for _ in presentation.slides]  # 슬라이드 수만큼 빈 리스트로 초기화

    # 프레젠테이션의 슬라이드를 순회
    for slide_num, slide in enumerate(presentation.slides):
        slide_charts = []

        # 슬라이드의 모든 쉐이프를 순회
        for shape in slide.shapes:
            # 차트 추출
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                chart = shape.chart
                chart_data = []

                # 차트 제목 추가
                if chart.has_title:
                    chart_data.append(f"Chart Title: {chart.chart_title.text_frame.text.strip()}")

                
                # 카테고리 데이터 추출
                categories = [category.label for category in chart.plots[0].categories]
                chart_data.append(f"Categories: {', '.join(categories)}")

                # 시리즈 데이터 추출
                for series in chart.plots[0].series:
                    series_name = series.name
                    series_values = [str(value) for value in series.values]
                    chart_data.append(f"  {series_name}: {series_values}")

                slide_charts.append("\n".join(chart_data))
        charts_data[slide_num] = slide_charts if slide_charts else ['']  # 슬라이드 인덱스에 차트 추가

    return charts_data


def extract_tables(ppt_file_path):
    # 파워포인트 프레젠테이션 파일 로드
    presentation = Presentation(ppt_file_path)

    # 표 데이터를 저장할 리스트
    tables_data = [[] for _ in presentation.slides]  # 슬라이드 수만큼 빈 리스트로 초기화

    # 프레젠테이션의 슬라이드를 순회
    for slide_num, slide in enumerate(presentation.slides):
        slide_tables = []

        # 슬라이드의 모든 쉐이프를 순회
        for shape in slide.shapes:
            # 표 추출
            if shape.has_table:
                table_data = []
                table = shape.table
                for row in table.rows:
                    table_data.append(" | ".join(cell.text.strip() for cell in row.cells))
                slide_tables.append("\n".join(table_data))

        tables_data[slide_num] = slide_tables if slide_tables else ['']  # 슬라이드 인덱스에 표 추가

    return tables_data