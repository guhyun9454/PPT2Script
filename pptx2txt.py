from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
import os

def extract_text_and_images(ppt_file_path):
    presentation = Presentation(ppt_file_path)
    base_name = os.path.splitext(os.path.basename(ppt_file_path))[0]
    text_result_dir = os.path.join(os.path.dirname(ppt_file_path), f'result_txt/{base_name}_text')
    
    # 텍스트 결과 디렉토리 생성
    os.makedirs(text_result_dir, exist_ok=True)

    text_content = []
    images_by_slide = {}

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = [f'--- Slide {slide_number} ---']
        images = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = paragraph.text.strip()
                    if paragraph_text:
                        slide_text.append(paragraph_text)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # 이미지를 메모리에 저장
                image_stream = io.BytesIO(shape.image.blob)
                image = Image.open(image_stream)
                images.append(image)
        
        if images:
            images_by_slide[f"Slide_{slide_number}"] = images
        
        if slide_text:
            text_content.append('\n'.join(slide_text))

    # 텍스트 파일로 결과 저장
    text_result_file = os.path.join(text_result_dir, f'{base_name}.txt')
    with open(text_result_file, 'w', encoding='utf-8') as file:
        file.write('\n\n'.join(text_content))
    
    print(f"Text extracted and saved to {text_result_file}")
    return images_by_slide

# 예시 사용 방법 (실제 사용시 주석 해제)
# extract_text_and_images('./test.pptx')


