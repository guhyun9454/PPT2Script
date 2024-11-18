from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

def extract_text_and_images(ppt_file_path):
    presentation = Presentation(ppt_file_path)
    text_content_by_slide = [[] for _ in presentation.slides]  # 슬라이드 수만큼 빈 문자열로 초기화
    images_by_slide = {}

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = []
        images = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = paragraph.text.strip()
                    if paragraph_text:
                        slide_text.append(paragraph_text)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # 이미지를 메모리에 저장
                image_stream = io.BytesIO(shape.image.blob)
                image = Image.open(image_stream)
                images.append(image)
        
        if images:
            images_by_slide[f"Slide {slide_number}"] = images
        
        if slide_text:
            text_content_by_slide[slide_number - 1] = '\n'.join(slide_text)  # 슬라이드 인덱스에 텍스트 추가

    return images_by_slide, text_content_by_slide

